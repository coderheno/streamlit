import streamlit as st
from PIL import Image
import io
import base64

st.set_page_config(page_title="My Education Journey", layout="wide")

st.title("My Journey through Education 1.0 to 5.0")
st.markdown("""
This interactive app is built to accompany your seminar introduction. Use the controls to navigate the five stages of your career and optionally upload photos for each stage.
""")

# Timeline data
stages = [
    {
        "stage": "Education 1.0",
        "years": "May 2010–May 2014",
        "title": "The Chalk and Blackboard Era",
        "place": "VLB Janakiammal Arts & Science College, Coimbatore — Assistant Professor",
        "text": "Taught mostly with chalk and board. One multimedia course with a rare projector attempt — the process felt like a big adventure! Technology in classrooms was minimal.",
        "default_image": None,
    },
    {
        "stage": "Education 2.0",
        "years": "Oct 2014–Oct 2019",
        "title": "Books and Projectors Take the Stage",
        "place": "Jigjiga University, Ethiopia",
        "text": "Depended on books and projectors. Faced challenges teaching courses like Data Structures but shone in Computer Graphics, adding visual and creative elements to lessons.",
        "default_image": None,
    },
    {
        "stage": "Education 3.0",
        "years": "Oct 2019–July 2023",
        "title": "The Internet and Interactive Learning Era",
        "place": "Skyline University, Nigeria",
        "text": "Technology became central. Used computers and the internet for learner-centric approaches, LO-based teaching, and rich multimedia resources to engage students.",
        "default_image": None,
    },
    {
        "stage": "Education 4.0",
        "years": "July 2023–Present",
        "title": "Smart Classrooms and Blended Learning",
        "place": "CHRIST University, India",
        "text": "Integrated smart classrooms and blended learning. Built ICT tools to replace PPTs, encouraged participatory learning, and made classes highly interactive.",
        "default_image": None,
    },
    {
        "stage": "Education 5.0",
        "years": "Future",
        "title": "A Magical, Student-Empowered Future, Dream Classroom",
        "place": "Right here",
        "text": "Students learn, create, and solve real problems using AI, VR, AR, and no-code tools. Technology fades into the background, but learning becomes limitless and magical.",
        "default_image": None,
    },
]

# Sidebar controls
st.sidebar.header("Navigation & Settings")
stage_index = st.sidebar.radio("Choose stage", list(range(1, len(stages) + 1)), format_func=lambda x: f"{x}. {stages[x-1]['stage']}") - 1
show_timeline = st.sidebar.checkbox("Show full timeline", value=True)
show_years = st.sidebar.checkbox("Show years on cards", value=True)



# Main area: timeline
if show_timeline:
    cols = st.columns(len(stages))
    for i, s in enumerate(stages):
        with cols[i]:
            st.markdown(f"**{s['stage']}**")
            if show_years:
                st.caption(s['years'])
            if s['default_image']:
                st.image(s['default_image'], use_column_width=True)
            else:
                # simple placeholder
                st.write("\n")
                st.info(s['title'])

st.markdown("---")

# Display selected stage details
s = stages[stage_index]
left, right = st.columns([2, 3])
with left:
    st.subheader(f"{s['stage']} — {s['years']}")
    st.markdown(f"**{s['title']}**")
    st.write(s['place'])
    st.write(s['text'])
    if s['default_image']:
        st.image(s['default_image'], caption="Uploaded photo", use_column_width=True)
   
with right:
    st.subheader("Quick points")
    if stage_index == 0:
        st.markdown("""
• Started career as assistant professor the Institute where i studied my Masters;  
• One multimedia attempt — big logistical effort.  
• Emphasis on fundamental teaching skills.
• Relied on chalk and blackboard. 
""")
        st.video("https://www.youtube.com/watch?v=R4elqE1I2RM")
    elif stage_index == 1:
        st.markdown("""
• Taught in Ethiopia; books and projectors were primary tools.  
• Faced challenges due to cultural, social backgrounds, Introduced visual approaches in Computer Graphics.
""")
        st.video("https://www.youtube.com/9eNZoPMDo3k")
        
    elif stage_index == 2:
        st.markdown("""
• Skyline University: integrated internet and LO-based methods,
    Focus on learner-centric pedagogy and digital resources.
""")
        st.video("https://www.youtube.com/watch?v=v5mjhe-lOSM")
    elif stage_index == 3:
        st.markdown("""
• Christ University: smart classrooms and blended learning.
Developed ICT tools and encouraged participatory learning.
""")
        st.video("https://youtu.be/8SUqGk0VgYE")
    
    else:
        st.markdown("""
• Vision for Education 5.0: student empowerment, tech-as-companion, focus on creativity and real-world problem solving.

""")
        
    

# Optional: export as PPTX
st.markdown("---")
st.subheader("Export")
if st.button("Generate simple PPTX of slides"):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        for s in stages:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"{s['stage']} — {s['years']}"
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, Inches(3))
            tf = txBox.text_frame
            tf.text = s['title'] + "\n\n" + s['place'] + "\n\n" + s['text']
            # add image if available
            if s['default_image']:
                img_io = io.BytesIO()
                s['default_image'].save(img_io, format='PNG')
                img_io.seek(0)
                slide.shapes.add_picture(img_io, Inches(6), Inches(1), height=Inches(3))

        bio = io.BytesIO()
        prs.save(bio)
        bio.seek(0)
        b64 = base64.b64encode(bio.read()).decode()
        href = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}"
        st.markdown(f"[Download PPTX]({href})", unsafe_allow_html=True)
    except Exception as e:
        st.error("Could not generate PPTX. Make sure python-pptx is installed in your environment.")
        st.write(e)
import streamlit as st

st.header("Session 2: Live Q&A & Polls with Slido")

st.write("Join the live interaction by clicking below:")

slido_url = "https://wall.sli.do/event/wFJqRfNF2Y2VmBR4LHwrSp?section=b9cfea71-b26b-464b-b6af-d6bf60a861d5"

# Display as a markdown link
st.markdown(f"[🔗 Join the Slido Q&A & Polls here →]({slido_url})")

# Or as a button
if st.button("Open Slido Session"):
    st.write(f"[Opening Slido... Click here if nothing happens]({slido_url})")
    
    st.markdown("""
**Quick Links:**  
[Session-2](https://13july2.streamlit.app)
""")
